from rest_framework.viewsets import ReadOnlyModelViewSet
from rest_framework.response import Response
from rest_framework.decorators import api_view
from rest_framework import status
from django.db.models import Sum, Q
from .models import VaccineData
from .serializers import VaccineDataSerializer
import json
from django.db.models.functions import TruncDate
from rest_framework.parsers import MultiPartParser, FormParser
from django.http import FileResponse
import csv
import io
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import base64

class VaccineDataViewSet(ReadOnlyModelViewSet):
    queryset = VaccineData.objects.all()
    serializer_class = VaccineDataSerializer
    filterset_fields = ["country", "state_or_region"]

@api_view(["GET"])
def get_comparison(request):
    """Retorna dados comparativos entre países"""
    countries = request.GET.getlist("countries", ["brasil", "portugal", "italia", "usa"])
    
    data = {}
    for country in countries:
        country_data = VaccineData.objects.filter(country=country).aggregate(
            total_vaccinated=Sum("vaccinated"),
            total_deaths=Sum("deaths"),
            total_population=Sum("population")
        )
        data[country] = country_data
    
    return Response(data)

@api_view(["GET"])
def get_chart_data(request):
    """Retorna dados para gráficos de evolução temporal"""
    chart_type = request.GET.get("type", "line")
    country = request.GET.get("country", "brasil")
    
    queryset = VaccineData.objects.filter(
        country=country
    ).values(
        "date", "vaccinated", "deaths"
    ).order_by("date").distinct()
    
    results = list(queryset)
    
    # Se não houver dados por data, agrupar por país
    if not results:
        results = [{
            "date": "Dados Agregados",
            "vaccinated": VaccineData.objects.filter(country=country).aggregate(Sum("vaccinated"))["vaccinated__sum"] or 0,
            "deaths": VaccineData.objects.filter(country=country).aggregate(Sum("deaths"))["deaths__sum"] or 0
        }]
    
    return Response(results)

@api_view(["GET"])
def get_countries_data(request):
    """Retorna totais por país com tratamento de dados vazios"""
    countries_list = request.GET.getlist("countries", ["brasil", "portugal", "italia", "usa"])
    
    results = []
    for country in countries_list:
        data = VaccineData.objects.filter(country=country).aggregate(
            vaccinated=Sum("vaccinated"),
            deaths=Sum("deaths")
        )
        results.append({
            "country": country,
            "vaccinated": data["vaccinated"] or 0,
            "deaths": data["deaths"] or 0
        })
    
    return Response(results)

@api_view(["GET"])
def get_state_data(request):
    """Retorna dados por estado/região com valores reais"""
    country = request.GET.get("country", "brasil")
    
    states = VaccineData.objects.filter(
        country=country
    ).values("state_or_region").distinct().order_by("state_or_region")
    
    results = []
    for state_obj in states:
        state = state_obj["state_or_region"]
        if not state:
            continue
            
        data = VaccineData.objects.filter(
            country=country,
            state_or_region=state
        ).aggregate(
            vaccinated=Sum("vaccinated"),
            deaths=Sum("deaths")
        )
        results.append({
            "state": state,
            "vaccinated": data["vaccinated"] or 0,
            "deaths": data["deaths"] or 0
        })
    
    # Ordenar por vacinados decrescente
    results.sort(key=lambda x: x["vaccinated"], reverse=True)
    
    return Response(results)

@api_view(["GET"])
def get_deaths_comparison(request):
    """Retorna comparação específica de óbitos entre países"""
    countries_list = request.GET.getlist("countries", ["brasil", "portugal", "italia", "usa"])
    
    results = []
    for country in countries_list:
        data = VaccineData.objects.filter(country=country).aggregate(
            total_deaths=Sum("deaths"),
            total_vaccinated=Sum("vaccinated")
        )
        total_deaths = data["total_deaths"] or 0
        total_vaccinated = data["total_vaccinated"] or 1  # Evitar divisão por zero
        
        # Taxa de mortalidade
        mortality_rate = (total_deaths / total_vaccinated * 100) if total_vaccinated > 0 else 0
        
        results.append({
            "country": country,
            "deaths": total_deaths,
            "vaccination_rate": total_vaccinated,
            "mortality_rate": round(mortality_rate, 2)
        })
    
    return Response(results)

@api_view(["POST"])
def upload_csv(request):
    """Upload de arquivo CSV para importar dados"""
    if 'file' not in request.FILES:
        return Response({"error": "Nenhum arquivo enviado"}, status=status.HTTP_400_BAD_REQUEST)
    
    csv_file = request.FILES['file']
    country_name = request.POST.get('country', 'custom_country')
    
    try:
        # Decodificar arquivo CSV
        decoded_file = csv_file.read().decode('utf-8')
        reader = csv.DictReader(io.StringIO(decoded_file))
        
        imported_count = 0
        for row in reader:
            try:
                # Esperado: date, state_or_region, vaccinated, deaths, population
                vaccine_data = VaccineData(
                    country=country_name.lower(),
                    state_or_region=row.get('state_or_region', 'Nacional'),
                    date=row['date'],
                    vaccinated=int(row.get('vaccinated', 0)),
                    deaths=int(row.get('deaths', 0)),
                    population=int(row.get('population', 0))
                )
                vaccine_data.save()
                imported_count += 1
            except Exception as e:
                continue
        
        return Response({
            "success": True,
            "message": f"{imported_count} registros importados com sucesso",
            "imported_count": imported_count
        })
    
    except Exception as e:
        return Response({
            "error": f"Erro ao processar arquivo: {str(e)}"
        }, status=status.HTTP_400_BAD_REQUEST)

@api_view(["GET"])
def export_powerpoint(request):
    """Gera apresentação PowerPoint com gráficos e análises"""
    try:
        # Criar apresentação
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Título
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(2)
        
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.text = "Análise de Vacinação"
        title_frame.paragraphs[0].font.size = Pt(60)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Comparativo Brasil, Portugal, Itália e EUA"
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Slide 2: Estatísticas Gerais
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Estatísticas Gerais"
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
        
        # Carregar dados
        countries_data = VaccineData.objects.values('country').distinct()
        countries = [c['country'] for c in countries_data]
        
        y_pos = 1.8
        for country in ['brasil', 'portugal', 'italia', 'usa']:
            country_stats = VaccineData.objects.filter(country=country).aggregate(
                vaccinated=Sum('vaccinated'),
                deaths=Sum('deaths')
            )
            
            text_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(9), Inches(0.9))
            text_frame = text_box.text_frame
            text = f"{country.upper()}: {country_stats['vaccinated'] or 0:,} vacinados | {country_stats['deaths'] or 0:,} óbitos"
            text_frame.text = text
            text_frame.paragraphs[0].font.size = Pt(16)
            y_pos += 1
        
        # Slide 3: Recomendações
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Análise e Recomendações"
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
        
        recommendations = [
            "Brasil: Manter esforços de vacinação em estados com baixa cobertura",
            "Portugal: Continuar monitoramento de casos em regiões de maior incidência",
            "Itália: Reforçar campanhas de vacinação em populações vulneráveis",
            "EUA: Intensificar ações de prevenção em áreas de baixa vacinação"
        ]
        
        y_pos = 1.8
        for rec in recommendations:
            text_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(8.6), Inches(0.8))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            text_frame.text = "• " + rec
            text_frame.paragraphs[0].font.size = Pt(14)
            y_pos += 1.0
        
        # Salvar arquivo
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        
        return FileResponse(
            pptx_io,
            as_attachment=True,
            filename=f"analise_vacinacao_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx",
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    
    except Exception as e:
        return Response({
            "error": f"Erro ao gerar PowerPoint: {str(e)}"
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)

@api_view(["GET"])
def export_csv(request):
    """Exporta todos os dados em CSV"""
    try:
        response = FileResponse(
            content_type='text/csv',
            as_attachment=True,
            filename=f"dados_vacinacao_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
        )
        
        writer = csv.writer(response)
        writer.writerow(['país', 'estado_região', 'data', 'vacinados', 'óbitos', 'população'])
        
        for data in VaccineData.objects.all():
            writer.writerow([
                data.country,
                data.state_or_region or 'N/A',
                data.date,
                data.vaccinated,
                data.deaths,
                data.population
            ])
        
        return response
    
    except Exception as e:
        return Response({
            "error": f"Erro ao exportar CSV: {str(e)}"
        }, status=status.HTTP_500_INTERNAL_SERVER_ERROR)
