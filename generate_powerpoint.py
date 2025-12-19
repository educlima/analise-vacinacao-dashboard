"""
Script para gerar apresentação PowerPoint com análise de dados
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import pandas as pd
from django.db import models
import os
from datetime import datetime

def generate_powerpoint_report(vaccine_data_queryset):
    """
    Gera relatório PowerPoint com análise completa de vacinação
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Cores da apresentação
    TEMA_AZUL = RGBColor(0, 102, 204)
    TEMA_CINZA = RGBColor(64, 64, 64)
    TEXTO_BRANCO = RGBColor(255, 255, 255)
    
    # Slide 1: Capa
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = TEMA_AZUL
    
    # Título
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "ANÁLISE DE DADOS DE VACINAÇÃO"
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = TEXTO_BRANCO
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"Brasil, Portugal, Itália e EUA\n{datetime.now().strftime('%B %Y')}"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = TEXTO_BRANCO
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 2: Resumo Executivo
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    title2.text = "Resumo Executivo"
    
    # Calcular totais
    data_df = pd.DataFrame(list(vaccine_data_queryset.values()))
    total_vacinados = data_df['vaccinated'].sum() if not data_df.empty else 0
    total_obitos = data_df['deaths'].sum() if not data_df.empty else 0
    total_paises = data_df['country'].nunique() if not data_df.empty else 0
    
    # Conteúdo
    content = slide2.placeholders[1].text_frame
    content.text = f"""Países Analisados: {total_paises}

Total de Pessoas Vacinadas: {total_vacinados:,.0f}

Total de Óbitos Registrados: {total_obitos:,.0f}

Período de Análise: Outubro 2025 - Dezembro 2025

Status: Análise em Tempo Real"""
    
    for paragraph in content.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.space_before = Pt(12)
    
    # Slide 3: Metodologia
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "Metodologia"
    
    content3 = slide3.placeholders[1].text_frame
    content3.text = """Fonte de Dados: Our World in Data e importações CSV customizadas

Métricas Utilizadas:
• Taxa de Vacinação = (Vacinados / População) × 100
• Taxa de Mortalidade = (Óbitos / População) × 100,000

Tratamento de Dados:
• Remoção de valores nulos
• Validação de tipos de dados
• Atualização periódica

Visualizações:
• Gráficos de barras para comparação
• Gráficos de linhas para tendências
• Gráficos de pizza para proporções"""
    
    for paragraph in content3.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.space_before = Pt(6)
    
    # Slide 4-7: Análise por País
    paises = ['brasil', 'portugal', 'italia', 'usa']
    nomes_paises = ['Brasil', 'Portugal', 'Itália', 'EUA']
    
    for idx, (pais, nome) in enumerate(zip(paises, nomes_paises)):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_slide = slide.shapes.title
        title_slide.text = f"Análise: {nome}"
        
        # Filtrar dados do país
        pais_data = data_df[data_df['country'] == pais]
        
        if not pais_data.empty:
            total_vac = pais_data['vaccinated'].sum()
            total_obt = pais_data['deaths'].sum()
            taxa_vac = pais_data['vaccination_rate'].mean()
            taxa_obt = pais_data['death_rate'].mean()
            
            content = slide.placeholders[1].text_frame
            content.text = f"""Vacinados: {total_vac:,.0f}

Óbitos: {total_obt:,.0f}

Taxa de Vacinação: {taxa_vac:.1f}%

Taxa de Mortalidade: {taxa_obt:.2f} por 100k hab

Análise: {nome} apresenta situação {'favorável' if taxa_vac > 80 else 'regular' if taxa_vac > 60 else 'crítica'} de vacinação."""
        else:
            content = slide.placeholders[1].text_frame
            content.text = f"Dados indisponíveis para {nome}"
        
        for paragraph in content.paragraphs:
            paragraph.font.size = Pt(18)
            paragraph.space_before = Pt(12)
    
    # Slide 8: Conclusões
    slide_final = prs.slides.add_slide(prs.slide_layouts[1])
    title_final = slide_final.shapes.title
    title_final.text = "Conclusões e Recomendações"
    
    content_final = slide_final.placeholders[1].text_frame
    content_final.text = """Principais Achados:
• Disparidades significativas entre países
• Correlação entre vacinação e redução de óbitos
• Necessidade de análise regional para otimização

Recomendações:
• Aumentar cobertura vacinal em regiões críticas
• Fortalecer sistema de vigilância epidemiológica
• Implementar campanhas direcionadas por região
• Continuar monitoramento de tendências

Próximas Ações:
• Inclusão de novos países na análise
• Análise de variantes e impacto na mortalidade
• Integração com dados socioeconômicos"""
    
    for paragraph in content_final.paragraphs:
        paragraph.font.size = Pt(16)
        paragraph.space_before = Pt(8)
    
    return prs

def save_powerpoint(presentation, filename='relatorio_vacinacao.pptx'):
    """Salva apresentação em arquivo"""
    presentation.save(filename)
    return filename
