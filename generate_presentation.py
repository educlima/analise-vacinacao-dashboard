"""
Script para gerar apresentação PowerPoint sobre o projeto
Documentação e arquitetura do sistema
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Cria apresentação completa do projeto"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Cores
    AZUL_PRINCIPAL = RGBColor(0, 102, 204)
    VERDE_DESTAQUE = RGBColor(0, 153, 76)
    CINZA_ESCURO = RGBColor(64, 64, 64)
    BRANCO = RGBColor(255, 255, 255)
    CINZA_CLARO = RGBColor(240, 240, 240)
    
    def add_title_slide(prs, titulo, subtitulo):
        """Cria slide de título"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = AZUL_PRINCIPAL
        
        # Título
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.text = titulo
        title_frame.paragraphs[0].font.size = Pt(54)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = BRANCO
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Subtítulo
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_frame.text = subtitulo
        subtitle_frame.paragraphs[0].font.size = Pt(28)
        subtitle_frame.paragraphs[0].font.color.rgb = BRANCO
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(prs, titulo, conteudo_list):
        """Cria slide de conteúdo com bullets"""
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = titulo
        
        # Formata título
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = AZUL_PRINCIPAL
        
        # Conteúdo
        content = slide.placeholders[1].text_frame
        content.clear()
        
        for idx, item in enumerate(conteudo_list):
            if idx == 0:
                p = content.paragraphs[0]
            else:
                p = content.add_paragraph()
            
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
            p.space_before = Pt(10)
            p.space_after = Pt(10)
        
        return slide
    
    # SLIDE 1: Capa
    add_title_slide(prs, "ANÁLISE DE DADOS DE VACINAÇÃO", 
                    "Sistema de Análise Comparativa\nBrasil, Portugal, Itália e EUA\nDezembro 2025")
    
    # SLIDE 2: Objetivo do Projeto
    add_content_slide(prs, "Objetivo do Projeto", [
        "Desenvolver sistema web para análise comparativa de dados de vacinação e óbitos",
        "Comparar estratégias de vacinação entre 4 países",
        "Identificar disparidades regionais e temporais",
        "Permitir importação de dados de novos países",
        "Gerar relatórios profissionais em PowerPoint"
    ])
    
    # SLIDE 3: Tecnologias Utilizadas
    add_content_slide(prs, "Stack Tecnológico", [
        "Backend: Django 5.0 + Django REST Framework",
        "Banco de Dados: PostgreSQL (Docker) / SQLite (desenvolvimento)",
        "Frontend: HTML5, CSS3, JavaScript, Plotly.js",
        "Containerização: Docker + Docker Compose",
        "Análise de Dados: Pandas, NumPy",
        "Relatórios: python-pptx, openpyxl"
    ])
    
    # SLIDE 4: Arquitetura do Sistema
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Arquitetura do Sistema"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = AZUL_PRINCIPAL
    
    content = slide.placeholders[1].text_frame
    content.text = """┌─────────────────────────────────────┐
│  NAVEGADOR (Frontend)              │
│  - Dashboard Interativo            │
│  - Gráficos Plotly                │
│  - Seletores Dinâmicos            │
└────────────┬────────────────────────┘
             │ HTTP/REST
┌────────────▼────────────────────────┐
│  DJANGO SERVER (Backend)           │
│  - API REST                        │
│  - Processamento de Dados          │
│  - Importação CSV                  │
└────────────┬────────────────────────┘
             │ SQL
┌────────────▼────────────────────────┐
│  POSTGRESQL DATABASE               │
│  - Dados de Vacinação             │
│  - Óbitos por região              │
└────────────────────────────────────┘"""
    
    for p in content.paragraphs:
        p.font.size = Pt(14)
        p.font.name = 'Courier New'
    
    # SLIDE 5: Funcionalidades Principais
    add_content_slide(prs, "Funcionalidades Principais", [
        "Dashboard interativo com gráficos em tempo real",
        "Múltiplos tipos de visualização (barras, linhas, pizza, dispersão)",
        "Filtros por país, estado/região e métrica",
        "Importação de arquivos CSV com novos países",
        "Exportação de dados em CSV e PowerPoint",
        "Análise comparativa de 4 países"
    ])
    
    # SLIDE 6: Fluxo de Importação CSV
    add_content_slide(prs, "Importação de Dados (CSV)", [
        "Usuário faz upload do arquivo CSV",
        "Sistema valida formato e estrutura dos dados",
        "Pandas processa e limpa os dados",
        "Cálculo automático de métricas (taxa vacinação, mortalidade)",
        "Dados salvos no banco com índices",
        "Dashboard atualizado com novo país"
    ])
    
    # SLIDE 7: Métricas Analisadas
    add_content_slide(prs, "Métricas e Indicadores", [
        "Total de Vacinados: Número absoluto de pessoas vacinadas",
        "Taxa de Vacinação: (Vacinados / População) × 100 %",
        "Total de Óbitos: Número absoluto de mortes registradas",
        "Taxa de Mortalidade: (Óbitos / População) × 100,000 hab",
        "Correlação: Relação entre vacinação e redução de óbitos"
    ])
    
    # SLIDE 8: APIs Disponíveis
    add_content_slide(prs, "APIs REST Implementadas", [
        "GET /api/countries-data/ → Dados agregados por país",
        "GET /api/state-data/ → Dados por estado/região",
        "GET /api/chart-data/ → Dados formatados para gráficos",
        "POST /api/import-csv/ → Importar novo arquivo CSV",
        "GET /api/export-csv/ → Exportar em CSV",
        "GET /api/generate-powerpoint/ → Gerar relatório"
    ])
    
    # SLIDE 9: Dashboard - Seletores
    add_content_slide(prs, "Dashboard: Seletores e Filtros", [
        "Tipo de Gráfico: Barras, Linhas, Pizza, Dispersão",
        "Métrica: Vacinados, Óbitos ou Ambos",
        "País: Todos, Brasil, Portugal, Itália, EUA (+ importados)",
        "Período: Seleção de data inicial e final",
        "Atualização em tempo real ao mudar filtros"
    ])
    
    # SLIDE 10: Gráficos Gerados
    add_content_slide(prs, "Gráficos Interativos", [
        "1. Comparativo entre Países: Confronto direto entre dados",
        "2. Evolução Temporal: Progressão ao longo do tempo",
        "3. Dados por Estado: Top 10 regiões",
        "4. Comparativo de Óbitos: Análise focada em mortalidade",
        "Todos com zoom, hover, download e legenda clicável"
    ])
    
    # SLIDE 11: Geração de PowerPoint
    add_content_slide(prs, "Exportação de Relatórios", [
        "Slide de Capa com título e data",
        "Resumo Executivo com totais e estatísticas",
        "Metodologia e fonte de dados",
        "Análise por país com gráficos",
        "Conclusões e recomendações",
        "Formato profissional pronto para apresentação"
    ])
    
    # SLIDE 12: Tratamento de Dados
    add_content_slide(prs, "Processamento de Dados", [
        "Validação de tipos (datas, números inteiros)",
        "Remoção de valores nulos e duplicados",
        "Normalização de nomes (países e regiões)",
        "Cálculo de métricas derivadas",
        "Indexação para buscas rápidas",
        "Integridade referencial entre registros"
    ])
    
    # SLIDE 13: Segurança
    add_content_slide(prs, "Medidas de Segurança", [
        "CSRF Protection contra ataques cross-site",
        "SQL Injection Prevention com queries parametrizadas",
        "XSS Prevention com template escaping",
        "Input Validation em uploads CSV",
        "CORS Headers para controle de origem",
        "Sem armazenamento de dados pessoais (dados agregados)"
    ])
    
    # SLIDE 14: Performance
    add_content_slide(prs, "Otimizações de Performance", [
        "Índices de banco de dados para buscas rápidas",
        "Caching de dados em memória",
        "Paginação de resultados (100 registros)",
        "Compressão Gzip de respostas HTTP",
        "Gráficos renderizados no cliente (Plotly.js)",
        "Tempo de resposta: < 500ms para maioria das operações"
    ])
    
    # SLIDE 15: Instalação (Windows)
    add_content_slide(prs, "Instalação: Windows 11", [
        "Pré-requisito: Docker Desktop instalado e rodando",
        "1. Extraia o arquivo ZIP no diretório desejado",
        "2. Abra terminal (PowerShell) na pasta do projeto",
        "3. Execute: docker-compose build --no-cache",
        "4. Execute: docker-compose up -d",
        "5. Aguarde 60 segundos e acesse: localhost:8000"
    ])
    
    # SLIDE 16: Instalação (macOS/Linux)
    add_content_slide(prs, "Instalação: macOS/Linux", [
        "Pré-requisito: Docker e Docker Compose instalados",
        "1. Clone ou extraia o projeto: cd vaccine-analysis",
        "2. Construir: docker-compose build --no-cache",
        "3. Iniciar: docker-compose up -d",
        "4. Migrar: docker-compose exec web python manage.py migrate",
        "5. Coletar dados: docker-compose exec web python scripts/collect_data.py",
        "6. Acessar: http://localhost:8000"
    ])
    
    # SLIDE 17: Estrutura de Diretórios
    add_content_slide(prs, "Estrutura do Projeto", [
        "config/ → Configurações Django (settings, urls, wsgi)",
        "vaccine/ → App principal (models, views, urls)",
        "templates/ → HTML para dashboard",
        "static/ → CSS e JavaScript",
        "scripts/ → Scripts Python (coleta de dados)",
        "docker-compose.yml → Orquestração de containers",
        "requirements.txt → Dependências Python"
    ])
    
    # SLIDE 18: Fluxo de Uso Típico
    add_content_slide(prs, "Fluxo de Uso Típico", [
        "1. Usuário acessa http://localhost:8000",
        "2. Dashboard carrega com gráficos dos 4 países",
        "3. Seleciona tipo gráfico, métrica e país",
        "4. Gráficos são regenerados em tempo real",
        "5. Importa arquivo CSV de novo país (opcional)",
        "6. Gera e baixa relatório PowerPoint",
        "7. Exporta dados em CSV para análise avançada"
    ])
    
    # SLIDE 19: Casos de Uso
    add_content_slide(prs, "Casos de Uso", [
        "Executivo: Visão geral com relatório PowerPoint",
        "Analista de Dados: Exporta CSV para análise em Excel",
        "Pesquisador: Estuda tendências com gráficos interativos",
        "Gestor de Saúde: Compara performance entre regiões",
        "Comunicação: Cria apresentações com dados visuais"
    ])
    
    # SLIDE 20: Problemas Comuns e Soluções
    add_content_slide(prs, "Troubleshooting Rápido", [
        "Docker não inicia → Verifique se Docker Desktop está rodando",
        "Porta 8000 em uso → Mude para porta diferente em docker-compose.yml",
        "Sem dados → Execute: docker-compose exec web python scripts/collect_data.py",
        "CSV não importa → Verifique colunas: date, state_or_region, vaccinated, deaths, population",
        "Dashboard lento → Reinicie: docker-compose restart web"
    ])
    
    # SLIDE 21: Próximas Funcionalidades
    add_content_slide(prs, "Roadmap Futuro", [
        "Autenticação de usuários com restrições de acesso",
        "Agendamento automático de coleta de dados",
        "Integração com APIs públicas de saúde em tempo real",
        "Análise preditiva com Machine Learning",
        "Dashboards personalizados por usuário",
        "Versão mobile responsiva",
        "Integração com Tableau ou Power BI"
    ])
    
    # SLIDE 22: Conclusão
    add_content_slide(prs, "Conclusão", [
        "Sistema completo de análise de dados de vacinação",
        "Interativo, escalável e fácil de usar",
        "Suporta importação de novos países",
        "Gera relatórios profissionais",
        "Pronto para uso em análises de saúde pública",
        "Código aberto e documentado"
    ])
    
    # SLIDE 23: Contatos e Repositório
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = AZUL_PRINCIPAL
    
    contact_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True
    contact_frame.text = """CONTATOS E RECURSOS

GitHub: github.com/seu-usuario/vaccine-analysis

Documentação: DOCUMENTACAO_TECNICA.md

Guia de Uso: GUIA_USO_DASHBOARD.md

Troubleshooting: TROUBLESHOOTING.md

Email: seu-email@exemplo.com"""
    
    for p in contact_frame.paragraphs:
        p.font.size = Pt(20)
        p.font.color.rgb = BRANCO
        p.space_before = Pt(6)
        p.alignment = PP_ALIGN.CENTER
    
    return prs

def save_presentation(filename='Analise de Vacinacao.pptx'):
    """Cria e salva a apresentação"""
    prs = create_presentation()
    prs.save(filename)
    return filename

if __name__ == '__main__':
    file_saved = save_presentation()
    print(f"Apresentação criada: {file_saved}")
